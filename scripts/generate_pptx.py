#!/usr/bin/env python3
"""Generate PowerPoint presentation for VNE research."""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

DIAGRAM_DIR = os.path.join(os.path.dirname(os.path.abspath(__file__)), "..", "docs", "diagrams")

# === Color Palette ===
BG_DARK = RGBColor(0x1B, 0x1B, 0x2F)       # Dark navy background
BG_CARD = RGBColor(0x24, 0x24, 0x3E)       # Card background
ACCENT_BLUE = RGBColor(0x4E, 0x9A, 0xF5)   # Primary accent
ACCENT_CYAN = RGBColor(0x56, 0xCC, 0xF2)   # Secondary accent
ACCENT_GREEN = RGBColor(0x48, 0xD1, 0xA5)  # Success/highlight
ACCENT_ORANGE = RGBColor(0xFF, 0x9F, 0x43) # Warning/emphasis
ACCENT_PURPLE = RGBColor(0xA2, 0x7A, 0xF0) # Purple accent
TEXT_WHITE = RGBColor(0xF0, 0xF0, 0xF0)    # Primary text
TEXT_GRAY = RGBColor(0xB0, 0xB0, 0xC0)     # Secondary text
TABLE_HEADER = RGBColor(0x2D, 0x3A, 0x5C)  # Table header bg
TABLE_ROW1 = RGBColor(0x1E, 0x1E, 0x36)    # Table row alt 1
TABLE_ROW2 = RGBColor(0x26, 0x26, 0x42)    # Table row alt 2


def set_slide_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height):
    return slide.shapes.add_textbox(left, top, width, height)


def set_text(tf, text, size=18, bold=False, color=TEXT_WHITE, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return p


def add_paragraph(tf, text, size=16, bold=False, color=TEXT_WHITE, alignment=PP_ALIGN.LEFT, space_before=Pt(6), space_after=Pt(2), font_name="Calibri"):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    p.space_before = space_before
    p.space_after = space_after
    return p


def add_bullet(tf, text, level=0, size=16, color=TEXT_WHITE, bold=False):
    p = tf.add_paragraph()
    p.text = text
    p.level = level
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = "Calibri"
    p.space_before = Pt(4)
    p.space_after = Pt(2)
    return p


def add_section_header(slide, number, title):
    """Add a section divider bar at top."""
    # Section number + title bar
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.9))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_BLUE
    shape.line.fill.background()
    tf = shape.text_frame
    tf.margin_left = Inches(0.5)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    set_text(tf, f"{number}. {title}", size=28, bold=True, color=TEXT_WHITE)


def add_slide_title(slide, title, subtitle=None):
    """Add slide title below section header."""
    tb = add_textbox(slide, Inches(0.6), Inches(1.0), Inches(12), Inches(0.6))
    set_text(tb.text_frame, title, size=26, bold=True, color=ACCENT_CYAN)
    if subtitle:
        tb2 = add_textbox(slide, Inches(0.6), Inches(1.55), Inches(12), Inches(0.4))
        set_text(tb2.text_frame, subtitle, size=16, color=TEXT_GRAY, bold=False)


def style_table(table, header_color=TABLE_HEADER, row1=TABLE_ROW1, row2=TABLE_ROW2):
    """Style a table with dark theme."""
    for row_idx, row in enumerate(table.rows):
        for cell in row.cells:
            cell.fill.solid()
            if row_idx == 0:
                cell.fill.fore_color.rgb = header_color
            elif row_idx % 2 == 1:
                cell.fill.fore_color.rgb = row1
            else:
                cell.fill.fore_color.rgb = row2
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(13)
                p.font.name = "Calibri"
                if row_idx == 0:
                    p.font.bold = True
                    p.font.color.rgb = ACCENT_CYAN
                else:
                    p.font.color.rgb = TEXT_WHITE
                p.alignment = PP_ALIGN.LEFT
            cell.margin_left = Inches(0.08)
            cell.margin_right = Inches(0.08)
            cell.margin_top = Inches(0.04)
            cell.margin_bottom = Inches(0.04)


def add_code_block(slide, left, top, width, height, code_text):
    """Add a styled code block."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x16, 0x16, 0x28)
    shape.line.color.rgb = RGBColor(0x3A, 0x3A, 0x5C)
    shape.line.width = Pt(1)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_top = Inches(0.15)
    lines = code_text.strip().split("\n")
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(11)
        p.font.name = "Consolas"
        p.font.color.rgb = ACCENT_GREEN
        p.space_before = Pt(1)
        p.space_after = Pt(1)


def add_card(slide, left, top, width, height, title, bullets, accent=ACCENT_BLUE):
    """Add a card-style box with title and bullets."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = BG_CARD
    shape.line.color.rgb = accent
    shape.line.width = Pt(1.5)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_top = Inches(0.1)
    set_text(tf, title, size=15, bold=True, color=accent)
    for b in bullets:
        add_bullet(tf, f"• {b}", size=12, color=TEXT_WHITE)


def add_equation_box(slide, left, top, width, height, equation):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x1A, 0x1A, 0x30)
    shape.line.color.rgb = ACCENT_PURPLE
    shape.line.width = Pt(1.5)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_top = Inches(0.08)
    set_text(tf, equation, size=14, color=ACCENT_CYAN, font_name="Consolas", alignment=PP_ALIGN.CENTER)


# ============================================================
# BUILD PRESENTATION
# ============================================================

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK_LAYOUT = prs.slide_layouts[6]

# ============================================================
# SLIDE 1: TITLE
# ============================================================
slide = prs.slides.add_slide(BLANK_LAYOUT)
set_slide_bg(slide)

# Title block
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(1.5), Inches(13.333), Inches(3.2))
shape.fill.solid()
shape.fill.fore_color.rgb = BG_CARD
shape.line.fill.background()

tb = add_textbox(slide, Inches(1), Inches(1.8), Inches(11.3), Inches(1.2))
set_text(tb.text_frame, "Virtual Network Embedding\nin Multi-Domain Networks", size=38, bold=True, color=ACCENT_CYAN, alignment=PP_ALIGN.CENTER)

tb = add_textbox(slide, Inches(1), Inches(3.2), Inches(11.3), Inches(0.8))
set_text(tb.text_frame, "A Study of PSO, Reinforcement Learning, and Swarm Intelligence Approaches", size=22, color=TEXT_GRAY, alignment=PP_ALIGN.CENTER)

tb = add_textbox(slide, Inches(1), Inches(4.2), Inches(11.3), Inches(0.5))
set_text(tb.text_frame, "Duvan Nguyen", size=20, bold=True, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)

# Accent line
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.5), Inches(3.05), Inches(4.3), Inches(0.04))
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT_BLUE
shape.line.fill.background()

# ============================================================
# SLIDE 2: OUTLINE
# ============================================================
slide = prs.slides.add_slide(BLANK_LAYOUT)
set_slide_bg(slide)
add_section_header(slide, "0", "OUTLINE")

sections = [
    ("I", "Abstract", ACCENT_BLUE),
    ("II", "Introduction & Motivation", ACCENT_CYAN),
    ("III", "Problem Formulation", ACCENT_GREEN),
    ("IV", "Related Work", ACCENT_ORANGE),
    ("V", "Proposed Methods", ACCENT_PURPLE),
    ("VI", "Experiment", ACCENT_BLUE),
    ("VII", "Results & Discussion", ACCENT_CYAN),
    ("VIII", "Conclusion & Future Work", ACCENT_GREEN),
    ("IX", "References", ACCENT_ORANGE),
]

for i, (num, title, color) in enumerate(sections):
    y = 1.2 + i * 0.62
    # Number box
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.5), Inches(y), Inches(0.8), Inches(0.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    set_text(tf, num, size=16, bold=True, color=TEXT_WHITE, alignment=PP_ALIGN.CENTER)
    # Title
    tb = add_textbox(slide, Inches(3.6), Inches(y), Inches(7), Inches(0.5))
    tf = tb.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    set_text(tf, title, size=20, color=TEXT_WHITE)

# ============================================================
# SLIDE 3: ABSTRACT
# ============================================================
slide = prs.slides.add_slide(BLANK_LAYOUT)
set_slide_bg(slide)
add_section_header(slide, "I", "ABSTRACT")

abstract_text = (
    "Network virtualization enables multiple virtual networks to coexist on shared physical "
    "infrastructure, but efficiently embedding virtual network requests onto substrate networks "
    "— the Virtual Network Embedding (VNE) problem — remains NP-hard, especially in multi-domain "
    "environments."
)
abstract_text2 = (
    "This work investigates and compares six VNE algorithms spanning three paradigms: "
    "heuristic (greedy), metaheuristic (Particle Swarm Optimization), and learning-based "
    "(Q-Learning, Deep Q-Networks, Swarm Reinforcement Learning). All algorithms operate "
    "within a unified centralized hierarchical multi-domain architecture."
)
abstract_text3 = (
    "Experiments on a multi-domain substrate network with 4 domains (120 nodes) show that "
    "hybrid approaches combining PSO with DQN-guided fitness and multi-path routing achieve "
    "the best trade-off between cost, delay, acceptance rate, and revenue-to-cost ratio. "
    "We contextualize our work against the Virne benchmark (ICLR 2026), which identifies "
    "generalization and scalability as key open challenges."
)

shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(1.2), Inches(11.3), Inches(5.5))
shape.fill.solid()
shape.fill.fore_color.rgb = BG_CARD
shape.line.color.rgb = ACCENT_BLUE
shape.line.width = Pt(2)
tf = shape.text_frame
tf.word_wrap = True
tf.margin_left = Inches(0.4)
tf.margin_top = Inches(0.3)
tf.margin_right = Inches(0.4)
set_text(tf, abstract_text, size=18, color=TEXT_WHITE)
add_paragraph(tf, "", size=10)
add_paragraph(tf, abstract_text2, size=18, color=TEXT_WHITE)
add_paragraph(tf, "", size=10)
add_paragraph(tf, abstract_text3, size=18, color=TEXT_WHITE)

# Keywords
add_paragraph(tf, "", size=8)
p = add_paragraph(tf, "Keywords: ", size=15, bold=True, color=ACCENT_CYAN)
keywords = "Virtual Network Embedding, Particle Swarm Optimization, Deep Q-Network, Swarm Intelligence, Multi-Domain, Multi-Path"
add_paragraph(tf, keywords, size=15, color=ACCENT_GREEN)

# ============================================================
# SLIDE 4: INTRODUCTION - What is Network Virtualization?
# ============================================================
slide = prs.slides.add_slide(BLANK_LAYOUT)
set_slide_bg(slide)
add_section_header(slide, "II", "INTRODUCTION & MOTIVATION")
add_slide_title(slide, "What is Network Virtualization?")

tb = add_textbox(slide, Inches(0.6), Inches(1.8), Inches(6), Inches(1))
tf = tb.text_frame
tf.word_wrap = True
set_text(tf, "Network virtualization decouples virtual networks from physical infrastructure, allowing multiple isolated virtual networks to share the same substrate.", size=16, color=TEXT_WHITE)

slide.shapes.add_picture(os.path.join(DIAGRAM_DIR, "network_virtualization.png"),
    Inches(0.3), Inches(2.8), width=Inches(6.2))

# Applications cards
add_card(slide, Inches(7), Inches(1.8), Inches(5.8), Inches(1.2),
    "Applications",
    ["Cloud computing: multi-tenant resource sharing",
     "5G network slicing: dedicated VN per service",
     "Internet of Drones (IoD) in Industry 4.0",
     "Edge computing for latency-sensitive services"],
    accent=ACCENT_GREEN)

add_card(slide, Inches(7), Inches(3.3), Inches(5.8), Inches(1.2),
    "VNE Problem (NP-hard)",
    ["Node mapping: assign virtual nodes → physical nodes",
     "Link mapping: route virtual links → physical paths",
     "Constraints: CPU capacity, bandwidth, domain"],
    accent=ACCENT_ORANGE)

add_card(slide, Inches(7), Inches(4.8), Inches(5.8), Inches(1.5),
    "Key Challenges in Multi-Domain",
    ["Combinatorial explosion: O(|Np|^|Nv|) possible mappings",
     "Coupled decisions: node placement affects link routing",
     "Cross-domain routing via boundary nodes",
     "Dynamic arrivals: online problem with resource fragmentation"],
    accent=ACCENT_PURPLE)

# ============================================================
# SLIDE 5: WHY IS VNE HARD / MOTIVATION
# ============================================================
slide = prs.slides.add_slide(BLANK_LAYOUT)
set_slide_bg(slide)
add_section_header(slide, "II", "INTRODUCTION & MOTIVATION")
add_slide_title(slide, "Motivation: Why Compare Multiple Approaches?")

tb = add_textbox(slide, Inches(0.6), Inches(1.8), Inches(12), Inches(0.8))
tf = tb.text_frame
tf.word_wrap = True
set_text(tf, "No single paradigm dominates VNE. We systematically compare three paradigms to understand their trade-offs:", size=17, color=TEXT_WHITE)

# Three paradigm cards
add_card(slide, Inches(0.6), Inches(2.8), Inches(3.8), Inches(2.5),
    "Heuristic",
    ["Greedy node ranking", "Kruskal MST for links", "Fast, simple, deterministic", "No global optimization"],
    accent=ACCENT_GREEN)

add_card(slide, Inches(4.7), Inches(2.8), Inches(3.8), Inches(2.5),
    "Metaheuristic",
    ["Particle Swarm Optimization", "Global search over solution space", "Genetic mutation avoids local optima", "No learning between requests"],
    accent=ACCENT_BLUE)

add_card(slide, Inches(8.8), Inches(2.8), Inches(3.8), Inches(2.5),
    "Learning-based",
    ["Q-Learning / Deep Q-Networks", "Swarm RL (multiple DQN agents)", "Adapts and improves over time", "Training cost, generalization issues"],
    accent=ACCENT_PURPLE)

# Arrow showing our hybrid
shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(5.6), Inches(10.3), Inches(1.2))
shape.fill.solid()
shape.fill.fore_color.rgb = BG_CARD
shape.line.color.rgb = ACCENT_ORANGE
shape.line.width = Pt(2)
tf = shape.text_frame
tf.word_wrap = True
tf.margin_left = Inches(0.3)
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
set_text(tf, "Our Approach: Combine all three → PSO + DQN Swarm + Multi-Path allocation", size=20, bold=True, color=ACCENT_ORANGE, alignment=PP_ALIGN.CENTER)
add_paragraph(tf, "6 algorithms from simple greedy to advanced Swarm RL + PSO hybrids, all in a unified multi-domain framework", size=15, color=TEXT_GRAY, alignment=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 6: PROBLEM FORMULATION - Substrate Network
# ============================================================
slide = prs.slides.add_slide(BLANK_LAYOUT)
set_slide_bg(slide)
add_section_header(slide, "III", "PROBLEM FORMULATION")
add_slide_title(slide, "Substrate Network Model", "Based on MP-VNE [Zhang et al., 2022]")

add_equation_box(slide, Inches(0.6), Inches(2.0), Inches(12), Inches(0.55),
    "Gˢ = ({Gᵢˢ}ᵢ₌₁ᴰ, Lˢ_inter)    |    Each domain: Gᵢˢ = (Nᵢˢ, Lᵢˢ)  with boundary nodes Bᵢˢ ⊂ Nᵢˢ")

# Node attributes table
tbl = slide.shapes.add_table(5, 4, Inches(0.6), Inches(2.9), Inches(5.8), Inches(2.2)).table
tbl.columns[0].width = Inches(1.2)
tbl.columns[1].width = Inches(1.0)
tbl.columns[2].width = Inches(2.0)
tbl.columns[3].width = Inches(1.6)
headers = ["Component", "Symbol", "Description", "Range"]
for i, h in enumerate(headers):
    tbl.cell(0, i).text = h
data = [
    ["Node CPU", "Cₙˢ", "Processing capacity", "[50, 100]"],
    ["Node Price", "Pₙˢ", "Unit price of CPU", "[1, 5]"],
    ["Node Delay", "Dₙˢ", "Processing delay", "[0.1, 2.0]"],
    ["Link BW", "Bₗˢ", "Bandwidth capacity", "[500, 1000]"],
]
for r, row in enumerate(data):
    for c, val in enumerate(row):
        tbl.cell(r+1, c).text = val
style_table(tbl)

# Inter-domain table
tbl2 = slide.shapes.add_table(4, 4, Inches(6.8), Inches(2.9), Inches(5.8), Inches(1.8)).table
tbl2.columns[0].width = Inches(1.4)
tbl2.columns[1].width = Inches(1.0)
tbl2.columns[2].width = Inches(1.8)
tbl2.columns[3].width = Inches(1.6)
headers2 = ["Inter-Domain", "Symbol", "Description", "Range"]
for i, h in enumerate(headers2):
    tbl2.cell(0, i).text = h
data2 = [
    ["Link BW", "Bₗˢ", "Bandwidth capacity", "[1000, 5000]"],
    ["Link Price", "Pₗˢ", "Unit price of BW", "[0.5, 2.0]"],
    ["Link Delay", "Dₗˢ", "Transmission delay", "[5.0, 20.0]"],
]
for r, row in enumerate(data2):
    for c, val in enumerate(row):
        tbl2.cell(r+1, c).text = val
style_table(tbl2)

tb = add_textbox(slide, Inches(0.6), Inches(5.3), Inches(12), Inches(0.5))
set_text(tb.text_frame, "Topology: Erdos-Renyi random graph, edge probability = 0.5 per domain, 4 domains × 30 nodes = 120 total nodes", size=15, color=TEXT_GRAY)

# ============================================================
# SLIDE 7: PROBLEM FORMULATION - Virtual Network
# ============================================================
slide = prs.slides.add_slide(BLANK_LAYOUT)
set_slide_bg(slide)
add_section_header(slide, "III", "PROBLEM FORMULATION")
add_slide_title(slide, "Virtual Network Request Model")

add_equation_box(slide, Inches(0.6), Inches(2.0), Inches(5.5), Inches(0.5),
    "Gᵛ = (Nᵛ, Lᵛ)")

# VN attributes
tbl = slide.shapes.add_table(4, 4, Inches(0.6), Inches(2.8), Inches(5.5), Inches(1.8)).table
tbl.columns[0].width = Inches(1.3)
tbl.columns[1].width = Inches(0.9)
tbl.columns[2].width = Inches(1.8)
tbl.columns[3].width = Inches(1.5)
for i, h in enumerate(["Component", "Symbol", "Description", "Range"]):
    tbl.cell(0, i).text = h
for r, row in enumerate([
    ["Node CPU", "Cᵢᵛ", "CPU demand", "[1, 20]"],
    ["Link BW", "Bᵢⱼᵛ", "Bandwidth demand", "[1, 50]"],
    ["Domain", "Dᵢᵛ", "Allowed domains", "Subset of {1..D}"],
]):
    for c, val in enumerate(row):
        tbl.cell(r+1, c).text = val
style_table(tbl)

# Arrival model
add_card(slide, Inches(6.8), Inches(2.0), Inches(5.8), Inches(2.8),
    "Request Arrival Model",
    ["Nodes per VNR: Uniform U(3, 10)",
     "Edge probability: 0.5",
     "Arrival process: Poisson, λ = 0.04",
     "Lifetime: Exponential, mean = 500 time units",
     "",
     "Dynamic online problem — requests arrive",
     "and depart continuously over simulation time"],
    accent=ACCENT_CYAN)

# ============================================================
# SLIDE 8: PROBLEM FORMULATION - Mapping & Constraints
# ============================================================
slide = prs.slides.add_slide(BLANK_LAYOUT)
set_slide_bg(slide)
add_section_header(slide, "III", "PROBLEM FORMULATION")
add_slide_title(slide, "VNE Mapping Definition & Constraints")

# Node mapping constraints
add_card(slide, Inches(0.6), Inches(2.0), Inches(5.8), Inches(2.5),
    "Node Mapping: Mₙ : Nᵛ → Nˢ",
    ["Resource: Cᵢᵛ ≤ C_avail(Mₙ(nᵢᵛ))",
     "Domain: Mₙ(nᵢᵛ) ∈ Domain(Dᵢᵛ)",
     "One-to-one: Mₙ(nᵢᵛ) ≠ Mₙ(nⱼᵛ)  ∀ i ≠ j",
     "",
     "Each virtual node maps to exactly one",
     "physical node with sufficient CPU"],
    accent=ACCENT_BLUE)

# Link mapping constraints
add_card(slide, Inches(6.8), Inches(2.0), Inches(5.8), Inches(2.5),
    "Link Mapping: Mₗ : Lᵛ → P(Lˢ)",
    ["Single-path: one physical path per virtual link",
     "Multi-path: split BW across ≤ 5 paths",
     "Bᵢⱼᵛ ≤ Σ B_avail(p)  for all paths p",
     "",
     "Multi-path improves acceptance rate by",
     "utilizing residual bandwidth across paths"],
    accent=ACCENT_GREEN)

# Objective
shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(4.8), Inches(12), Inches(1.8))
shape.fill.solid()
shape.fill.fore_color.rgb = BG_CARD
shape.line.color.rgb = ACCENT_ORANGE
shape.line.width = Pt(2)
tf = shape.text_frame
tf.word_wrap = True
tf.margin_left = Inches(0.3)
tf.margin_top = Inches(0.15)
set_text(tf, "Optimization Objective", size=18, bold=True, color=ACCENT_ORANGE)
add_paragraph(tf, "min Cost(Gᵛ) = Σ Cᵛₙ × Pˢ_M(nᵛ)  +  Σ Bᵛₗ × Pˢₗ × hops", size=16, color=ACCENT_CYAN, font_name="Consolas", alignment=PP_ALIGN.CENTER)
add_paragraph(tf, "                    Node cost                    Link cost", size=13, color=TEXT_GRAY, font_name="Consolas", alignment=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 9: PROBLEM FORMULATION - Metrics
# ============================================================
slide = prs.slides.add_slide(BLANK_LAYOUT)
set_slide_bg(slide)
add_section_header(slide, "III", "PROBLEM FORMULATION")
add_slide_title(slide, "Evaluation Metrics")

tbl = slide.shapes.add_table(6, 3, Inches(1), Inches(2.0), Inches(11), Inches(3.5)).table
tbl.columns[0].width = Inches(1.8)
tbl.columns[1].width = Inches(5.0)
tbl.columns[2].width = Inches(4.2)
for i, h in enumerate(["Metric", "Formula", "Description"]):
    tbl.cell(0, i).text = h
metrics = [
    ["RAC", "# Accepted VNRs / # Total VNRs", "Request Acceptance Rate"],
    ["LAR", "Σ REV(v) × lifetime(v) / T", "Long-term Average Revenue"],
    ["R2C", "Σ REV × lifetime / Σ COST × lifetime", "Revenue-to-Cost Ratio"],
    ["Avg Cost", "Σ COST(v) / # Successes", "Average per-embedding cost"],
    ["Avg Delay", "Σ PathDelay(v) / # Successes", "Average physical path latency"],
]
for r, row in enumerate(metrics):
    for c, val in enumerate(row):
        tbl.cell(r+1, c).text = val
style_table(tbl)

tb = add_textbox(slide, Inches(1), Inches(5.7), Inches(11), Inches(0.6))
set_text(tb.text_frame, "Revenue proxy: REV(v) = Σ Cᵛₙ + Σ Bᵛₗ  (total resource demand)", size=15, color=TEXT_GRAY)

# ============================================================
# SLIDE 10: RELATED WORK - Overview
# ============================================================
slide = prs.slides.add_slide(BLANK_LAYOUT)
set_slide_bg(slide)
add_section_header(slide, "IV", "RELATED WORK")
add_slide_title(slide, "Three Paradigms for Solving VNE")

tbl = slide.shapes.add_table(4, 4, Inches(0.6), Inches(2.2), Inches(12), Inches(3.0)).table
tbl.columns[0].width = Inches(2.0)
tbl.columns[1].width = Inches(3.5)
tbl.columns[2].width = Inches(3.5)
tbl.columns[3].width = Inches(3.0)
for i, h in enumerate(["Paradigm", "Approach", "Strengths", "Weaknesses"]):
    tbl.cell(0, i).text = h
rows = [
    ["Heuristic", "Greedy ranking, MST", "Fast, simple, deterministic", "No global optimization, myopic"],
    ["Metaheuristic", "PSO, GA, ACO", "Global search, avoids local optima", "Slow, no learning between requests"],
    ["Learning-based", "Q-Learning, DQN, GNN+RL", "Adapts over time, learns patterns", "Training cost, generalization"],
]
for r, row in enumerate(rows):
    for c, val in enumerate(row):
        tbl.cell(r+1, c).text = val
style_table(tbl)

tb = add_textbox(slide, Inches(0.6), Inches(5.5), Inches(12), Inches(1))
tf = tb.text_frame
tf.word_wrap = True
set_text(tf, "Our work spans all three paradigms with 6 algorithms, exploring how combining them yields better results.", size=17, color=TEXT_GRAY)

# ============================================================
# SLIDE 11: RELATED WORK - MP-VNE
# ============================================================
slide = prs.slides.add_slide(BLANK_LAYOUT)
set_slide_bg(slide)
add_section_header(slide, "IV", "RELATED WORK")
add_slide_title(slide, "MP-VNE [Zhang et al., 2022]", "Multi-Domain VNE Based on Multi-Objective Optimization for IoD")

add_card(slide, Inches(0.6), Inches(2.0), Inches(5.8), Inches(3.5),
    "Key Contributions",
    ["Centralized hierarchical multi-domain architecture",
     "  with Global-Local controllers",
     "PSO with genetic mutation (10% random reset)",
     "  to avoid local optima",
     "Estimated mapping cost for candidate",
     "  node pre-selection",
     "Weighted summation converting multi-objective",
     "  (cost + delay) into single objective"],
    accent=ACCENT_BLUE)

add_card(slide, Inches(6.8), Inches(2.0), Inches(5.8), Inches(1.8),
    "Results",
    ["Outperforms MC-VNM, VNE-PSO, LID-VNE",
     "Acceptance rate: ~60% (steady)",
     "Mapping cost: 650-750 (best)",
     "Mapping delay: ~460 (minimum)"],
    accent=ACCENT_GREEN)

add_card(slide, Inches(6.8), Inches(4.1), Inches(5.8), Inches(1.4),
    "Our Adoption",
    ["Problem formulation and multi-domain architecture",
     "PSO framework as foundation for all 6 algorithms",
     "Multi-path link allocation mechanism"],
    accent=ACCENT_ORANGE)

# ============================================================
# SLIDE 12: RELATED WORK - FlagVNE
# ============================================================
slide = prs.slides.add_slide(BLANK_LAYOUT)
set_slide_bg(slide)
add_section_header(slide, "IV", "RELATED WORK")
add_slide_title(slide, "FlagVNE [Wang et al., IJCAI 2024]", "Flexible and Generalizable RL Framework for Network Resource Allocation")

add_card(slide, Inches(0.6), Inches(2.0), Inches(3.8), Inches(3.8),
    "Bidirectional Action",
    ["Joint selection of virtual AND",
     "physical node simultaneously",
     "aₜ = (nᵛ, nᵖ)",
     "",
     "Expands action space from",
     "|Nᵖ|×1 to |Nᵖ|×|Nᵛ|",
     "",
     "Provably better than",
     "unidirectional (Theorem 1)"],
    accent=ACCENT_BLUE)

add_card(slide, Inches(4.7), Inches(2.0), Inches(3.8), Inches(3.8),
    "Hierarchical Decoder",
    ["Bilevel policy decomposition:",
     "π(aₜ|sₜ) = πᴴ(nᵛ|sₜ) · πᴸ(nᵖ|sₜ,nᵛ)",
     "",
     "High-level: which virtual node",
     "Low-level: where to place it",
     "",
     "Reduces distribution size:",
     "|Nᵛ|×|Nᵖ| → |Nᵛ|+|Nᵖ|"],
    accent=ACCENT_CYAN)

add_card(slide, Inches(8.8), Inches(2.0), Inches(3.8), Inches(3.8),
    "Meta-RL + Curriculum",
    ["MAML for varying VNR sizes",
     "as distinct tasks",
     "",
     "Curriculum scheduling:",
     "1. Start with smallest VNRs",
     "2. Monitor policy entropy",
     "3. Add larger sizes when",
     "   entropy < threshold δ",
     "",
     "Result: +10.4% RAC over A3C-GCN"],
    accent=ACCENT_GREEN)

add_card(slide, Inches(0.6), Inches(6.1), Inches(12), Inches(0.8),
    "Our Reference",
    ["FlagVNE represents state-of-the-art RL-VNE. Bidirectional actions and meta-RL are future work directions for our system."],
    accent=ACCENT_ORANGE)

# ============================================================
# SLIDE 13: RELATED WORK - Swarm RL
# ============================================================
slide = prs.slides.add_slide(BLANK_LAYOUT)
set_slide_bg(slide)
add_section_header(slide, "IV", "RELATED WORK")
add_slide_title(slide, "Swarm RL Using PSO [Srini-Rohan, GitHub]")

add_card(slide, Inches(0.6), Inches(2.0), Inches(5.8), Inches(2.2),
    "Core Idea",
    ["Multiple DQN agents act as PSO particles",
     "Each agent explores independently, shares discoveries",
     "PSO equations influence Q-value updates:",
     "  Q_target += β(Q_pBest - Q_curr) + δ(Q_gBest - Q_curr)",
     "Personal best (pBest) + Global best (gBest) tracking"],
    accent=ACCENT_PURPLE)

add_card(slide, Inches(6.8), Inches(2.0), Inches(5.8), Inches(2.2),
    "Results (CartPole Benchmark)",
    ["Single DQN agent: 130 episodes to converge",
     "4-agent PSO swarm: 117 episodes (best agent)",
     "~10% faster convergence through collaboration",
     "",
     "Demonstrates: multi-agent knowledge sharing",
     "accelerates RL convergence"],
    accent=ACCENT_GREEN)

add_card(slide, Inches(0.6), Inches(4.5), Inches(12), Inches(1.5),
    "Our Adoption: Adapted Swarm RL for VNE",
    ["SRL-VNE: DQN swarm for node selection + single-path Dijkstra",
     "MP-DQN-VNE: DQN-guided PSO fitness + multi-path allocation",
     "SRL-MP-VNE: Full swarm RL + PSO + multi-path (most advanced hybrid)"],
    accent=ACCENT_ORANGE)

# ============================================================
# SLIDE 14: RELATED WORK - Virne
# ============================================================
slide = prs.slides.add_slide(BLANK_LAYOUT)
set_slide_bg(slide)
add_section_header(slide, "IV", "RELATED WORK")
add_slide_title(slide, "Virne Benchmark [Wang et al., ICLR 2026]", "Comprehensive Benchmark for Deep RL-based NFV Resource Allocation")

tbl = slide.shapes.add_table(6, 3, Inches(0.6), Inches(2.2), Inches(7), Inches(3.0)).table
tbl.columns[0].width = Inches(2.2)
tbl.columns[1].width = Inches(2.4)
tbl.columns[2].width = Inches(2.4)
for i, h in enumerate(["Feature", "Previous Benchmarks", "Virne"]):
    tbl.cell(0, i).text = h
rows = [
    ["Scenarios", "Cloud only", "Cloud + Edge + 5G"],
    ["Algorithms", "1-5", "30+"],
    ["RL support", "None", "Gym-style environments"],
    ["Evaluation", "Effectiveness only", "Multi-perspective"],
    ["GNN support", "None", "MLP/CNN/GCN/GAT/DualGCN"],
]
for r, row in enumerate(rows):
    for c, val in enumerate(row):
        tbl.cell(r+1, c).text = val
style_table(tbl)

add_card(slide, Inches(8), Inches(2.2), Inches(4.6), Inches(3.0),
    "Key Findings from Virne",
    ["PPO-DualGAT best: 78.1% RAC, 0.74 R2C",
     "Fixed intermediate reward (0.1) is optimal",
     "Action masking: +5.3% acceptance",
     "Generalization is biggest open problem",
     "  → policies degrade on unseen traffic",
     "",
     "We use Virne's metrics & methodology"],
    accent=ACCENT_CYAN)

# ============================================================
# SLIDE 15: PROPOSED METHODS - Architecture
# ============================================================
slide = prs.slides.add_slide(BLANK_LAYOUT)
set_slide_bg(slide)
add_section_header(slide, "V", "PROPOSED METHODS")
add_slide_title(slide, "System Architecture: Centralized Hierarchical Multi-Domain")

slide.shapes.add_picture(os.path.join(DIAGRAM_DIR, "system_architecture.png"),
    Inches(0.3), Inches(1.8), width=Inches(8.0))

add_card(slide, Inches(8.5), Inches(2.0), Inches(4.2), Inches(2.0),
    "Global Controller",
    ["Orchestrates optimization algorithm",
     "Inter-domain path routing",
     "Commit/rollback allocations"],
    accent=ACCENT_BLUE)

add_card(slide, Inches(8.5), Inches(4.3), Inches(4.2), Inches(2.0),
    "Local Controllers",
    ["Per-domain resource management",
     "Candidate node selection",
     "Intra-domain path caching"],
    accent=ACCENT_GREEN)

# ============================================================
# SLIDE 16: PROPOSED METHODS - Two Phase Pipeline
# ============================================================
slide = prs.slides.add_slide(BLANK_LAYOUT)
set_slide_bg(slide)
add_section_header(slide, "V", "PROPOSED METHODS")
add_slide_title(slide, "Two-Phase Mapping Pipeline")

slide.shapes.add_picture(os.path.join(DIAGRAM_DIR, "two_phase_pipeline.png"),
    Inches(0.3), Inches(1.8), width=Inches(6.8))

add_card(slide, Inches(7.5), Inches(2.0), Inches(5.2), Inches(1.8),
    "Single-Path Link Mapping",
    ["One physical path per virtual link",
     "Dijkstra or Kruskal MST",
     "Simple but may fail if path lacks BW",
     "Used by: MC-VNM, MPQ-VNE, SRL-VNE"],
    accent=ACCENT_CYAN)

add_card(slide, Inches(7.5), Inches(4.1), Inches(5.2), Inches(2.2),
    "Multi-Path Link Mapping (≤ 5 paths)",
    ["Split bandwidth across multiple paths",
     "1. Find shortest path → allocate min(demand, avail)",
     "2. Update remaining demand",
     "3. Repeat until satisfied or 5 paths",
     "4. Rollback if insufficient",
     "Used by: MP-VNE, MP-DQN-VNE, SRL-MP-VNE"],
    accent=ACCENT_GREEN)

# ============================================================
# SLIDE 17: PROPOSED METHODS - Algorithm Summary Table
# ============================================================
slide = prs.slides.add_slide(BLANK_LAYOUT)
set_slide_bg(slide)
add_section_header(slide, "V", "PROPOSED METHODS")
add_slide_title(slide, "Algorithm Summary")

tbl = slide.shapes.add_table(7, 5, Inches(0.6), Inches(2.0), Inches(12), Inches(3.5)).table
tbl.columns[0].width = Inches(2.0)
tbl.columns[1].width = Inches(3.0)
tbl.columns[2].width = Inches(3.0)
tbl.columns[3].width = Inches(2.0)
tbl.columns[4].width = Inches(2.0)
for i, h in enumerate(["Algorithm", "Node Mapping", "Link Mapping", "Learning", "Paths"]):
    tbl.cell(0, i).text = h
algos = [
    ["MC-VNM", "Greedy ranking", "Kruskal MST", "None", "Single"],
    ["MP-VNE", "PSO (20p, 15i)", "Floyd shortest path", "None", "Multi (≤5)"],
    ["MPQ-VNE", "Q-Learning (ε-greedy)", "Dijkstra", "Q-table", "Single"],
    ["SRL-VNE", "DQN + PSO Swarm", "Dijkstra", "DQN (4 agents)", "Single"],
    ["MP-DQN-VNE", "PSO + DQN fitness", "Floyd shortest path", "DQN (4 agents)", "Multi (≤5)"],
    ["SRL-MP-VNE", "PSO + Swarm RL", "Floyd shortest path", "DQN (4 agents)", "Multi (≤5)"],
]
for r, row in enumerate(algos):
    for c, val in enumerate(row):
        tbl.cell(r+1, c).text = val
style_table(tbl)

# Evolution arrow
slide.shapes.add_picture(os.path.join(DIAGRAM_DIR, "algorithm_evolution.png"),
    Inches(0.3), Inches(5.4), width=Inches(12.5))

# ============================================================
# SLIDE 18: METHOD 1 - MC-VNM
# ============================================================
slide = prs.slides.add_slide(BLANK_LAYOUT)
set_slide_bg(slide)
add_section_header(slide, "V", "PROPOSED METHODS")
add_slide_title(slide, "Method 1: MC-VNM — Baseline Greedy")

add_card(slide, Inches(0.6), Inches(2.0), Inches(5.8), Inches(2.8),
    "Node Mapping — Greedy Cost-Efficient Selection",
    ["1. Sort virtual nodes by CPU demand (descending)",
     "   → hardest-to-place first",
     "2. Score each candidate substrate node:",
     "   Score(nˢ) = P_cpu(nˢ) / (C_avail(nˢ) + ε)",
     "3. Select node with LOWEST score",
     "   → cheapest per available CPU unit"],
    accent=ACCENT_BLUE)

add_card(slide, Inches(6.8), Inches(2.0), Inches(5.8), Inches(2.8),
    "Link Mapping — Kruskal's MST",
    ["1. Sort substrate links by cost:",
     "   D_trans(l) + P_bw(l) × B_demand",
     "2. Build MST using Union-Find",
     "3. Find path in MST via BFS",
     "4. Single path per virtual link"],
    accent=ACCENT_GREEN)

add_card(slide, Inches(0.6), Inches(5.1), Inches(5.8), Inches(1.5),
    "Pros",
    ["Fastest algorithm — no iterations, no learning",
     "Deterministic — reproducible results",
     "Simple to implement and debug"],
    accent=ACCENT_GREEN)

add_card(slide, Inches(6.8), Inches(5.1), Inches(5.8), Inches(1.5),
    "Cons",
    ["Myopic — each node placed independently",
     "No global optimization of mapping cost",
     "Single-path limits bandwidth flexibility"],
    accent=ACCENT_ORANGE)

# ============================================================
# SLIDE 19: METHOD 2 - MP-VNE
# ============================================================
slide = prs.slides.add_slide(BLANK_LAYOUT)
set_slide_bg(slide)
add_section_header(slide, "V", "PROPOSED METHODS")
add_slide_title(slide, "Method 2: MP-VNE — PSO + Multi-Path")

add_card(slide, Inches(0.6), Inches(2.0), Inches(7.5), Inches(1.8),
    "Particle Swarm Optimization for Node Mapping",
    ["Each particle = a complete node mapping scheme",
     "Swarm of particles searches solution space collectively",
     "Velocity: vᵢ = w·vᵢ + c₁·r₁·(pbestᵢ - xᵢ) + c₂·r₂·(gbest - xᵢ)",
     "Position: xᵢ = xᵢ + vᵢ"],
    accent=ACCENT_BLUE)

# PSO params table
tbl = slide.shapes.add_table(7, 3, Inches(0.6), Inches(4.1), Inches(5.5), Inches(2.8)).table
tbl.columns[0].width = Inches(1.8)
tbl.columns[1].width = Inches(1.2)
tbl.columns[2].width = Inches(2.5)
for i, h in enumerate(["Parameter", "Value", "Role"]):
    tbl.cell(0, i).text = h
params = [
    ["Particles", "20", "Population size"],
    ["Iterations", "15", "Search depth"],
    ["w (inertia)", "0.7", "Momentum of direction"],
    ["c₁ (cognitive)", "1.5", "Pull toward personal best"],
    ["c₂ (social)", "1.5", "Pull toward global best"],
    ["Mutation rate", "0.1", "10% random reset"],
]
for r, row in enumerate(params):
    for c, val in enumerate(row):
        tbl.cell(r+1, c).text = val
style_table(tbl)

add_card(slide, Inches(6.5), Inches(4.1), Inches(6.2), Inches(1.2),
    "Fitness Function",
    ["f(x) = Σ NodeCost + Σ LinkCost",
     "LinkCost = D_trans + P_bw × B_demand per hop"],
    accent=ACCENT_CYAN)

add_card(slide, Inches(6.5), Inches(5.6), Inches(6.2), Inches(1.3),
    "Genetic Mutation (Key Innovation)",
    ["10% probability of random position reset",
     "Prevents PSO from converging to local optima",
     "Inspired by genetic algorithm crossover"],
    accent=ACCENT_ORANGE)

# ============================================================
# SLIDE 19b: MP-VNE — Multi-Path Allocation Diagram
# ============================================================
slide = prs.slides.add_slide(BLANK_LAYOUT)
set_slide_bg(slide)
add_section_header(slide, "V", "PROPOSED METHODS")
add_slide_title(slide, "Method 2: MP-VNE — Multi-Path Link Allocation")

slide.shapes.add_picture(os.path.join(DIAGRAM_DIR, "multipath_allocation.png"),
    Inches(1.5), Inches(1.8), width=Inches(10.3))

# ============================================================
# SLIDE 20: METHOD 3 - MPQ-VNE
# ============================================================
slide = prs.slides.add_slide(BLANK_LAYOUT)
set_slide_bg(slide)
add_section_header(slide, "V", "PROPOSED METHODS")
add_slide_title(slide, "Method 3: MPQ-VNE — Q-Learning")

add_card(slide, Inches(0.6), Inches(2.0), Inches(5.8), Inches(3.5),
    "Q-Learning Node Selection",
    ["Q-table maps (domain_id, snode_id) → Q-value",
     "",
     "Initialization:",
     "  Q(s) = C_cpu(nˢ) / (P_cpu(nˢ) + ε)",
     "",
     "Action selection (ε-greedy):",
     "  90%: select argmax Q(s)  (exploit)",
     "  10%: select random  (explore)",
     "",
     "Q-update after each VNR:",
     "  Q(s) ← Q(s) + α·(reward - Q(s))"],
    accent=ACCENT_BLUE)

tbl = slide.shapes.add_table(5, 2, Inches(6.8), Inches(2.0), Inches(5.5), Inches(2.0)).table
tbl.columns[0].width = Inches(3.0)
tbl.columns[1].width = Inches(2.5)
for i, h in enumerate(["Parameter", "Value"]):
    tbl.cell(0, i).text = h
for r, row in enumerate([
    ["Learning rate (α)", "0.1"],
    ["Discount factor (γ)", "0.9"],
    ["Exploration (ε)", "0.1"],
    ["Reward (success/fail)", "+1.0 / -1.0"],
]):
    for c, val in enumerate(row):
        tbl.cell(r+1, c).text = val
style_table(tbl)

add_card(slide, Inches(6.8), Inches(4.3), Inches(5.5), Inches(1.1),
    "Link Mapping",
    ["Dijkstra's single shortest path",
     "Cost: D_trans + P_bw × B_demand"],
    accent=ACCENT_GREEN)

add_card(slide, Inches(6.8), Inches(5.7), Inches(5.5), Inches(1.0),
    "Limitation",
    ["Tabular Q-learning doesn't scale to large",
     "or continuous state spaces → need DQN"],
    accent=ACCENT_ORANGE)

# ============================================================
# SLIDE 21: METHOD 4 - SRL-VNE
# ============================================================
slide = prs.slides.add_slide(BLANK_LAYOUT)
set_slide_bg(slide)
add_section_header(slide, "V", "PROPOSED METHODS")
add_slide_title(slide, "Method 4: SRL-VNE — DQN + PSO Swarm")

add_card(slide, Inches(0.6), Inches(2.0), Inches(4.0), Inches(2.0),
    "DQN Architecture",
    ["Input: 6-dim state vector",
     "→ FC 64 → ReLU",
     "→ FC 64 → ReLU",
     "→ Output: 20 Q-values"],
    accent=ACCENT_BLUE)

# State table
tbl = slide.shapes.add_table(7, 3, Inches(4.9), Inches(2.0), Inches(4.5), Inches(2.8)).table
tbl.columns[0].width = Inches(1.2)
tbl.columns[1].width = Inches(2.0)
tbl.columns[2].width = Inches(1.3)
for i, h in enumerate(["Feature", "Description", "Norm"]):
    tbl.cell(0, i).text = h
feats = [
    ["v_cpu", "Virtual node CPU demand", "÷ 20"],
    ["s_cpu", "Substrate node avail CPU", "÷ 100"],
    ["s_bw", "Avg available bandwidth", "÷ 500"],
    ["s_degree", "Node degree in network", "÷ 1000"],
    ["progress", "Mapping completion ratio", "[0,1]"],
    ["bias", "Constant term", "0.5"],
]
for r, row in enumerate(feats):
    for c, val in enumerate(row):
        tbl.cell(r+1, c).text = val
style_table(tbl)

add_card(slide, Inches(9.7), Inches(2.0), Inches(3.0), Inches(2.8),
    "Swarm (4 Agents)",
    ["Each agent = DQN",
     "Track pBest, gBest",
     "",
     "PSO Q-update:",
     "Q += β(Q_pBest-Q)",
     "   + δ(Q_gBest-Q)",
     "",
     "β = δ = 0.1"],
    accent=ACCENT_PURPLE)

add_card(slide, Inches(0.6), Inches(5.1), Inches(5.5), Inches(1.6),
    "Training Details",
    ["Memory buffer: 2000 transitions, batch: 64",
     "ε: 1.0 → 0.05 (decay 0.995 per episode)",
     "γ (discount): 0.95"],
    accent=ACCENT_CYAN)

add_card(slide, Inches(6.5), Inches(5.1), Inches(6.1), Inches(1.6),
    "Reward Function",
    ["R = 100 × Revenue / (Cost + ε)",
     "Revenue = Σ CPU_demand + Σ BW_demand",
     "Link mapping: Dijkstra (single path)"],
    accent=ACCENT_GREEN)

# ============================================================
# SLIDE 21b: SRL-VNE — Swarm PSO Diagram
# ============================================================
slide = prs.slides.add_slide(BLANK_LAYOUT)
set_slide_bg(slide)
add_section_header(slide, "V", "PROPOSED METHODS")
add_slide_title(slide, "Method 4: SRL-VNE — PSO Swarm Knowledge Sharing")

slide.shapes.add_picture(os.path.join(DIAGRAM_DIR, "pso_swarm.png"),
    Inches(1.2), Inches(1.8), width=Inches(10.8))

# ============================================================
# SLIDE 22: METHOD 5 - MP-DQN-VNE
# ============================================================
slide = prs.slides.add_slide(BLANK_LAYOUT)
set_slide_bg(slide)
add_section_header(slide, "V", "PROPOSED METHODS")
add_slide_title(slide, "Method 5: MP-DQN-VNE — PSO Guided by DQN")

add_card(slide, Inches(0.6), Inches(2.0), Inches(12), Inches(0.9),
    "Key Innovation: DQN acts as an intelligent fitness function for PSO",
    ["PSO handles combinatorial search; DQN provides learned heuristics about node quality. DQN improves over time → PSO fitness becomes more accurate."],
    accent=ACCENT_ORANGE)

slide.shapes.add_picture(os.path.join(DIAGRAM_DIR, "mp_dqn_pipeline.png"),
    Inches(0.2), Inches(3.0), width=Inches(7.8))

add_card(slide, Inches(8.5), Inches(3.2), Inches(4.2), Inches(1.4),
    "Link Mapping",
    ["Multi-path allocation (up to 5 paths)",
     "Same as MP-VNE mechanism",
     "Floyd shortest path routing"],
    accent=ACCENT_GREEN)

add_card(slide, Inches(8.5), Inches(4.9), Inches(4.2), Inches(1.3),
    "Training",
    ["All 4 swarm agents receive same reward",
     "Cooperative learning",
     "R = 100 × Revenue / (Cost + ε)"],
    accent=ACCENT_PURPLE)

# ============================================================
# SLIDE 23: METHOD 6 - SRL-MP-VNE
# ============================================================
slide = prs.slides.add_slide(BLANK_LAYOUT)
set_slide_bg(slide)
add_section_header(slide, "V", "PROPOSED METHODS")
add_slide_title(slide, "Method 6: SRL-MP-VNE — Full Swarm RL + Multi-Path", "Most complete hybrid — combines all techniques")

tbl = slide.shapes.add_table(6, 2, Inches(0.6), Inches(2.2), Inches(5.0), Inches(2.5)).table
tbl.columns[0].width = Inches(2.5)
tbl.columns[1].width = Inches(2.5)
for i, h in enumerate(["Component", "Configuration"]):
    tbl.cell(0, i).text = h
for r, row in enumerate([
    ["PSO particles", "10"],
    ["PSO iterations", "50 (deepest search)"],
    ["Swarm agents", "4 DQN agents"],
    ["Max paths/link", "5"],
    ["Reward", "+100 (success) / -50 (fail)"],
]):
    for c, val in enumerate(row):
        tbl.cell(r+1, c).text = val
style_table(tbl)

slide.shapes.add_picture(os.path.join(DIAGRAM_DIR, "srl_mp_pipeline.png"),
    Inches(5.8), Inches(2.0), width=Inches(6.8))

# Comparison with MP-DQN-VNE
tbl2 = slide.shapes.add_table(4, 3, Inches(0.6), Inches(5.0), Inches(12), Inches(1.8)).table
tbl2.columns[0].width = Inches(4.0)
tbl2.columns[1].width = Inches(4.0)
tbl2.columns[2].width = Inches(4.0)
for i, h in enumerate(["", "MP-DQN-VNE", "SRL-MP-VNE"]):
    tbl2.cell(0, i).text = h
for r, row in enumerate([
    ["PSO iterations", "10", "50 (deeper search)"],
    ["Reward signal", "Continuous R2C", "Binary (+100/-50)"],
    ["Degree normalization", "÷ 1000", "÷ 10"],
]):
    for c, val in enumerate(row):
        tbl2.cell(r+1, c).text = val
style_table(tbl2)

# ============================================================
# SLIDE 24: EXPERIMENT SETUP - Substrate
# ============================================================
slide = prs.slides.add_slide(BLANK_LAYOUT)
set_slide_bg(slide)
add_section_header(slide, "VI", "EXPERIMENT")
add_slide_title(slide, "Experimental Setup — Substrate Network")

tbl = slide.shapes.add_table(7, 2, Inches(0.6), Inches(2.0), Inches(5.5), Inches(3.0)).table
tbl.columns[0].width = Inches(3.0)
tbl.columns[1].width = Inches(2.5)
for i, h in enumerate(["Parameter", "Value"]):
    tbl.cell(0, i).text = h
for r, row in enumerate([
    ["Number of domains", "4"],
    ["Nodes per domain", "30"],
    ["Boundary nodes/domain", "2"],
    ["Total physical nodes", "120"],
    ["Intra-domain edge prob", "0.5"],
    ["Inter-domain edge prob", "0.1"],
]):
    for c, val in enumerate(row):
        tbl.cell(r+1, c).text = val
style_table(tbl)

tbl2 = slide.shapes.add_table(7, 3, Inches(6.5), Inches(2.0), Inches(6.2), Inches(3.0)).table
tbl2.columns[0].width = Inches(2.2)
tbl2.columns[1].width = Inches(2.0)
tbl2.columns[2].width = Inches(2.0)
for i, h in enumerate(["Resource", "Intra-domain", "Inter-domain"]):
    tbl2.cell(0, i).text = h
for r, row in enumerate([
    ["CPU capacity", "U(50, 100)", "—"],
    ["CPU price", "U(1, 5)", "—"],
    ["Processing delay", "U(0.1, 2.0)", "—"],
    ["BW capacity", "U(500, 1000)", "U(1000, 5000)"],
    ["BW price", "U(0.1, 1.0)", "U(0.5, 2.0)"],
    ["Trans. delay", "U(1.0, 10.0)", "U(5.0, 20.0)"],
]):
    for c, val in enumerate(row):
        tbl2.cell(r+1, c).text = val
style_table(tbl2)

# ============================================================
# SLIDE 25: EXPERIMENT SETUP - VNR & Evaluation
# ============================================================
slide = prs.slides.add_slide(BLANK_LAYOUT)
set_slide_bg(slide)
add_section_header(slide, "VI", "EXPERIMENT")
add_slide_title(slide, "Experimental Setup — Virtual Requests & Evaluation Protocol")

tbl = slide.shapes.add_table(7, 2, Inches(0.6), Inches(2.0), Inches(5.5), Inches(2.8)).table
tbl.columns[0].width = Inches(3.0)
tbl.columns[1].width = Inches(2.5)
for i, h in enumerate(["VNR Parameter", "Value"]):
    tbl.cell(0, i).text = h
for r, row in enumerate([
    ["Nodes per VNR", "U(3, 10)"],
    ["Edge probability", "0.5"],
    ["CPU demand", "U(1, 20)"],
    ["Bandwidth demand", "U(1, 50)"],
    ["Arrival process", "Poisson, λ=0.04"],
    ["Lifetime", "Exp, mean=500"],
]):
    for c, val in enumerate(row):
        tbl.cell(r+1, c).text = val
style_table(tbl)

add_card(slide, Inches(6.5), Inches(2.0), Inches(6.2), Inches(2.0),
    "Evaluation Protocol",
    ["3 independent runs per algorithm, same dataset",
     "Results averaged with ± std deviation bands",
     "6 metrics plotted over simulation time",
     "Binned every 1000 time units"],
    accent=ACCENT_BLUE)

tbl3 = slide.shapes.add_table(7, 2, Inches(6.5), Inches(4.3), Inches(6.2), Inches(2.8)).table
tbl3.columns[0].width = Inches(1.8)
tbl3.columns[1].width = Inches(4.4)
for i, h in enumerate(["Metric", "What it Measures"]):
    tbl3.cell(0, i).text = h
for r, row in enumerate([
    ["RAC", "Request acceptance rate over time"],
    ["LAR", "Long-term average revenue"],
    ["R2C", "Revenue-to-cost efficiency"],
    ["Avg Cost", "Average per-embedding cost"],
    ["Avg Delay", "Average physical path latency"],
    ["Success #", "Cumulative successful embeddings"],
]):
    for c, val in enumerate(row):
        tbl3.cell(r+1, c).text = val
style_table(tbl3)

# ============================================================
# SLIDE 26: RESULTS - Plots
# ============================================================
slide = prs.slides.add_slide(BLANK_LAYOUT)
set_slide_bg(slide)
add_section_header(slide, "VII", "RESULTS & DISCUSSION")
add_slide_title(slide, "Performance Comparison — 6 Metrics Over Time")

# Placeholder for chart image
shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(2.0), Inches(11.7), Inches(4.8))
shape.fill.solid()
shape.fill.fore_color.rgb = BG_CARD
shape.line.color.rgb = ACCENT_BLUE
shape.line.width = Pt(2)
tf = shape.text_frame
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
set_text(tf, "[Insert: results/test_1/algorithm_comparison_plots.png]", size=24, color=TEXT_GRAY, alignment=PP_ALIGN.CENTER)
add_paragraph(tf, "6 subplots: RAC, LAR, R2C, Avg Cost, Avg Delay, Success Count", size=16, color=TEXT_GRAY, alignment=PP_ALIGN.CENTER)
add_paragraph(tf, "Each line = one algorithm, shaded bands = ± std deviation across 3 runs", size=14, color=TEXT_GRAY, alignment=PP_ALIGN.CENTER)

# Try to insert image
import os
img_path = "/Users/duvannguyen/Workspace/Studies/Virtual-Network-Embedding/results/test_1/algorithm_comparison_plots.png"
if os.path.exists(img_path):
    # Remove the placeholder shape
    sp = shape._element
    sp.getparent().remove(sp)
    slide.shapes.add_picture(img_path, Inches(0.5), Inches(1.8), width=Inches(12.3))

# ============================================================
# SLIDE 27: RESULTS - Analysis
# ============================================================
slide = prs.slides.add_slide(BLANK_LAYOUT)
set_slide_bg(slide)
add_section_header(slide, "VII", "RESULTS & DISCUSSION")
add_slide_title(slide, "Analysis")

tbl = slide.shapes.add_table(6, 2, Inches(0.6), Inches(2.0), Inches(12), Inches(2.8)).table
tbl.columns[0].width = Inches(3.0)
tbl.columns[1].width = Inches(9.0)
for i, h in enumerate(["Aspect", "Finding"]):
    tbl.cell(0, i).text = h
for r, row in enumerate([
    ["Acceptance Rate", "Multi-path methods (MP-VNE, MP-DQN-VNE, SRL-MP-VNE) achieve higher RAC — flexible BW splitting"],
    ["Embedding Cost", "PSO-based methods find lower-cost mappings than greedy MC-VNM"],
    ["Learning Effect", "RL methods show improving performance over time as DQN agents learn"],
    ["Delay", "PSO fitness explicitly includes delay → lower average latency than greedy"],
    ["Revenue", "Higher acceptance + lower cost → higher long-term revenue for hybrid methods"],
]):
    for c, val in enumerate(row):
        tbl.cell(r+1, c).text = val
style_table(tbl)

# Speed-quality tradeoff
tbl2 = slide.shapes.add_table(7, 5, Inches(0.6), Inches(5.2), Inches(12), Inches(2.0)).table
tbl2.columns[0].width = Inches(2.2)
tbl2.columns[1].width = Inches(2.2)
tbl2.columns[2].width = Inches(2.6)
tbl2.columns[3].width = Inches(2.5)
tbl2.columns[4].width = Inches(2.5)
for i, h in enumerate(["Algorithm", "Speed", "Quality", "Scalability", "Learning"]):
    tbl2.cell(0, i).text = h
for r, row in enumerate([
    ["MC-VNM", "★★★★★ Fastest", "★☆ Lowest", "★★★★ Good", "None"],
    ["MP-VNE", "★★★ Medium", "★★★ Good", "★★★ Medium", "None"],
    ["MPQ-VNE", "★★★★ Fast", "★★ Medium", "★★★★ Good", "Q-table"],
    ["SRL-VNE", "★★ Slower", "★★★ Good", "★★★ Medium", "DQN"],
    ["MP-DQN-VNE", "★ Slowest", "★★★★ High", "★★★ Medium", "DQN+PSO"],
    ["SRL-MP-VNE", "★ Slowest", "★★★★★ Highest", "★★★ Medium", "Swarm"],
]):
    for c, val in enumerate(row):
        tbl2.cell(r+1, c).text = val
style_table(tbl2)

# ============================================================
# SLIDE 28: RESULTS - Comparison with Literature
# ============================================================
slide = prs.slides.add_slide(BLANK_LAYOUT)
set_slide_bg(slide)
add_section_header(slide, "VII", "RESULTS & DISCUSSION")
add_slide_title(slide, "Comparison with Literature")

tbl = slide.shapes.add_table(5, 4, Inches(0.6), Inches(2.0), Inches(12), Inches(2.5)).table
tbl.columns[0].width = Inches(3.5)
tbl.columns[1].width = Inches(3.0)
tbl.columns[2].width = Inches(2.5)
tbl.columns[3].width = Inches(3.0)
for i, h in enumerate(["Source", "Best Method", "RAC", "Key Insight"]):
    tbl.cell(0, i).text = h
for r, row in enumerate([
    ["MP-VNE [Zhang 2022]", "MP-VNE", "~60%", "PSO + mutation for multi-domain"],
    ["FlagVNE [IJCAI 2024]", "FlagVNE", "+10.4% vs A3C-GCN", "Bidirectional + meta-RL"],
    ["Virne [ICLR 2026]", "PPO-DualGAT", "78.1%", "GNN encoder + action masking"],
    ["Our work", "SRL-MP-VNE", "(see plots)", "PSO + Swarm DQN + multi-path"],
]):
    for c, val in enumerate(row):
        tbl.cell(r+1, c).text = val
style_table(tbl)

add_card(slide, Inches(0.6), Inches(4.8), Inches(5.8), Inches(2.0),
    "Our Unique Contributions",
    ["Multi-domain focus (vs single-domain in Virne/FlagVNE)",
     "Systematic comparison of 3 paradigms in unified framework",
     "Novel PSO + DQN hybrid combinations",
     "Multi-path link allocation across all PSO methods"],
    accent=ACCENT_BLUE)

add_card(slide, Inches(6.8), Inches(4.8), Inches(5.8), Inches(2.0),
    "Gaps to Address",
    ["No GNN encoder yet (Virne: DualGAT is best)",
     "No bidirectional actions (FlagVNE advantage)",
     "No meta-RL for generalization",
     "Limited to synthetic topology (120 nodes)"],
    accent=ACCENT_ORANGE)

# ============================================================
# SLIDE 29: CONCLUSION
# ============================================================
slide = prs.slides.add_slide(BLANK_LAYOUT)
set_slide_bg(slide)
add_section_header(slide, "VIII", "CONCLUSION & FUTURE WORK")
add_slide_title(slide, "Conclusion")

add_card(slide, Inches(0.6), Inches(2.0), Inches(5.8), Inches(2.5),
    "What We Built",
    ["Unified multi-domain VNE framework with",
     "  Global-Local controller architecture",
     "6 algorithms: greedy → PSO → Q-Learning",
     "  → DQN Swarm → PSO+DQN → Full hybrid",
     "Multi-path link allocation (≤ 5 paths)",
     "Automated evaluation (6 metrics, 3 runs)"],
    accent=ACCENT_BLUE)

add_card(slide, Inches(6.8), Inches(2.0), Inches(5.8), Inches(2.5),
    "Key Findings",
    ["PSO >> Greedy (global search matters)",
     "DQN-guided PSO combines learned heuristics",
     "  with metaheuristic search effectively",
     "Swarm intelligence (4 agents + PSO sharing)",
     "  accelerates RL convergence",
     "Multi-path ↑ acceptance, ↑ complexity",
     "Clear speed–quality trade-off exists"],
    accent=ACCENT_GREEN)

# ============================================================
# SLIDE 30: FUTURE WORK
# ============================================================
slide = prs.slides.add_slide(BLANK_LAYOUT)
set_slide_bg(slide)
add_section_header(slide, "VIII", "CONCLUSION & FUTURE WORK")
add_slide_title(slide, "Future Work")

tbl = slide.shapes.add_table(7, 3, Inches(0.6), Inches(2.0), Inches(12), Inches(4.5)).table
tbl.columns[0].width = Inches(3.2)
tbl.columns[1].width = Inches(5.8)
tbl.columns[2].width = Inches(3.0)
for i, h in enumerate(["Direction", "Description", "Reference"]):
    tbl.cell(0, i).text = h
futures = [
    ["GNN-based encoders", "GCN, GAT, DualGAT to capture graph topology in state representation", "Virne [ICLR 2026]"],
    ["Bidirectional actions", "Joint virtual-physical node selection — provably better search space", "FlagVNE [IJCAI 2024]"],
    ["Meta-RL + curriculum", "Generalize across VNR sizes without retraining from scratch", "FlagVNE [IJCAI 2024]"],
    ["Latency / energy aware", "Multi-objective optimization for edge and 5G scenarios", "Virne [ICLR 2026]"],
    ["Larger-scale evaluation", "Real topologies: GEANT (40), BRAIN, WX500 (500 nodes)", "Virne [ICLR 2026]"],
    ["Virne integration", "Standardized comparison against 30+ algorithms in unified benchmark", "Virne [ICLR 2026]"],
]
for r, row in enumerate(futures):
    for c, val in enumerate(row):
        tbl.cell(r+1, c).text = val
style_table(tbl)

# ============================================================
# SLIDE 31: REFERENCES
# ============================================================
slide = prs.slides.add_slide(BLANK_LAYOUT)
set_slide_bg(slide)
add_section_header(slide, "IX", "REFERENCES")

refs = [
    '[1] P. Zhang, C. Wang, Z. Qin, H. Cao, "A Multi-Domain VNE Algorithm Based on Multi-Objective Optimization for IoD Architecture in Industry 4.0," arXiv:2202.12830, 2022.',
    '[2] T. Wang, Q. Fan, C. Wang, L. Yang, L. Ding, N. J. Yuan, H. Xiong, "FlagVNE: A Flexible and Generalizable Reinforcement Learning Framework for Network Resource Allocation," IJCAI, 2024.',
    '[3] Srini-Rohan, "Swarm Reinforcement Learning Using PSO," GitHub.',
    '[4] T. Wang, L. Deng, X. Chen, J. Wang, H. He, L. Ding, W. Wu, Q. Fan, H. Xiong, "Virne: A Comprehensive Benchmark for Deep RL-based Network Resource Allocation in NFV," ICLR, 2026.',
]

tb = add_textbox(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(5.5))
tf = tb.text_frame
tf.word_wrap = True
set_text(tf, refs[0], size=16, color=TEXT_WHITE)
for ref in refs[1:]:
    add_paragraph(tf, "", size=10)
    add_paragraph(tf, ref, size=16, color=TEXT_WHITE)

# ============================================================
# SLIDE 32: APPENDIX - Path Finding
# ============================================================
slide = prs.slides.add_slide(BLANK_LAYOUT)
set_slide_bg(slide)
add_section_header(slide, "A", "APPENDIX")
add_slide_title(slide, "Path Finding Algorithms Used")

tbl = slide.shapes.add_table(4, 4, Inches(0.6), Inches(2.0), Inches(12), Inches(2.0)).table
tbl.columns[0].width = Inches(3.0)
tbl.columns[1].width = Inches(4.5)
tbl.columns[2].width = Inches(2.5)
tbl.columns[3].width = Inches(2.0)
for i, h in enumerate(["Algorithm", "Used By", "Type", "Complexity"]):
    tbl.cell(0, i).text = h
for r, row in enumerate([
    ["Floyd-Warshall", "MP-VNE, MP-DQN-VNE, SRL-MP-VNE", "All-pairs (cached)", "O(V³)"],
    ["Dijkstra", "MPQ-VNE, SRL-VNE", "Single-source", "O(E log V)"],
    ["Kruskal MST+BFS", "MC-VNM", "Min spanning tree", "O(E log E)"],
]):
    for c, val in enumerate(row):
        tbl.cell(r+1, c).text = val
style_table(tbl)

add_equation_box(slide, Inches(0.6), Inches(4.3), Inches(12), Inches(0.6),
    "Link cost function:  Cost(lˢ) = D_transmission(lˢ) + P_bandwidth(lˢ) × B_demand(lᵛ)")

# ============================================================
# SLIDE 33: APPENDIX - DQN Hyperparameters
# ============================================================
slide = prs.slides.add_slide(BLANK_LAYOUT)
set_slide_bg(slide)
add_section_header(slide, "A", "APPENDIX")
add_slide_title(slide, "DQN Hyperparameter Comparison")

tbl = slide.shapes.add_table(10, 4, Inches(1), Inches(2.0), Inches(11), Inches(4.5)).table
tbl.columns[0].width = Inches(3.0)
tbl.columns[1].width = Inches(2.5)
tbl.columns[2].width = Inches(2.75)
tbl.columns[3].width = Inches(2.75)
for i, h in enumerate(["Parameter", "SRL-VNE", "MP-DQN-VNE", "SRL-MP-VNE"]):
    tbl.cell(0, i).text = h
params = [
    ["Swarm agents", "4", "4", "4"],
    ["Network architecture", "6→64→64→20", "6→64→64→20", "6→64→64→1"],
    ["Memory buffer", "2000", "2000", "2000"],
    ["Batch size", "64", "64", "64"],
    ["γ (discount)", "0.95", "0.95", "0.95"],
    ["ε start → end", "1.0 → 0.05", "1.0 → 0.05", "1.0 → 0.05"],
    ["ε decay rate", "0.995", "0.995", "0.995"],
    ["β (pBest influence)", "0.1", "0.1", "0.1"],
    ["δ (gBest influence)", "0.1", "0.1", "0.1"],
]
for r, row in enumerate(params):
    for c, val in enumerate(row):
        tbl.cell(r+1, c).text = val
style_table(tbl)

# ============================================================
# SLIDE 34: THANK YOU
# ============================================================
slide = prs.slides.add_slide(BLANK_LAYOUT)
set_slide_bg(slide)

shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(2.0), Inches(13.333), Inches(3.5))
shape.fill.solid()
shape.fill.fore_color.rgb = BG_CARD
shape.line.fill.background()

tb = add_textbox(slide, Inches(1), Inches(2.5), Inches(11.3), Inches(1.0))
set_text(tb.text_frame, "Thank You", size=48, bold=True, color=ACCENT_CYAN, alignment=PP_ALIGN.CENTER)

tb = add_textbox(slide, Inches(1), Inches(3.6), Inches(11.3), Inches(0.6))
set_text(tb.text_frame, "Questions?", size=28, color=TEXT_GRAY, alignment=PP_ALIGN.CENTER)

shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.5), Inches(3.35), Inches(4.3), Inches(0.04))
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT_BLUE
shape.line.fill.background()

# ============================================================
# SAVE
# ============================================================
output_path = "/Users/duvannguyen/Workspace/Studies/Virtual-Network-Embedding/docs/VNE_Presentation.pptx"
prs.save(output_path)
print(f"Presentation saved to {output_path}")
print(f"Total slides: {len(prs.slides)}")
